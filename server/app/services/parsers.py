import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import pytesseract
from app.config import settings

try:
    from google.cloud import vision
except Exception:
    vision = None

def parse_pdf(path: str) -> str:
    text_chunks = []
    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_chunks.append(t)
    except Exception:
        pass
    return "\n".join(text_chunks)

def parse_docx(path: str) -> str:
    try:
        doc = DocxDocument(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception:
        return ""

def parse_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""

def parse_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""

def parse_image(path: str) -> str:
    # Prefer Cloud Vision if enabled and available
    if settings.use_vision and vision is not None:
        try:
            client = vision.ImageAnnotatorClient()
            with open(path, "rb") as image_file:
                content = image_file.read()
            image = vision.Image(content=content)
            response = client.text_detection(image=image)
            if response.error.message:
                # fallback to tesseract on error
                raise RuntimeError(response.error.message)
            annotations = response.text_annotations
            if annotations:
                return annotations[0].description
        except Exception:
            # Fallback to Tesseract
            pass
    try:
        img = Image.open(path)
        return pytesseract.image_to_string(img)
    except Exception:
        return ""
